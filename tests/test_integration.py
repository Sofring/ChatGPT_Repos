from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE

from svg_to_pptx.cli import main

DATA_DIR = Path(__file__).parent / "data"


def _convert(svg_name: str, tmp_path) -> Path:
    svg_path = DATA_DIR / svg_name
    output_path = tmp_path / f"{svg_path.stem}.pptx"
    main([str(svg_path), str(output_path)])
    return output_path


def test_simple_shapes_conversion(tmp_path):
    pptx_path = _convert("simple_shapes.svg", tmp_path)
    prs = Presentation(pptx_path)
    shapes = list(prs.slides[0].shapes)
    assert len(shapes) == 3

    rect = next(shape for shape in shapes if shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.RECTANGLE)
    assert rect.fill.fore_color.rgb == RGBColor(0xFF, 0x00, 0x00)

    circle = next(shape for shape in shapes if shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.OVAL and shape.width == shape.height)
    assert circle.fill.fore_color.rgb == RGBColor(0x00, 0xFF, 0x00)

    line = next(shape for shape in shapes if shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.LINE)
    assert line.line.color.rgb == RGBColor(0x00, 0x00, 0xFF)


def test_text_conversion(tmp_path):
    pptx_path = _convert("text.svg", tmp_path)
    prs = Presentation(pptx_path)
    shapes = list(prs.slides[0].shapes)
    assert len(shapes) == 1

    textbox = shapes[0]
    assert textbox.has_text_frame
    assert textbox.text_frame.text == "Hello SVG Text"

    run = textbox.text_frame.paragraphs[0].runs[0]
    assert run.font.color.rgb == RGBColor(0x33, 0x33, 0x33)


def test_path_conversion(tmp_path):
    pptx_path = _convert("path.svg", tmp_path)
    prs = Presentation(pptx_path)
    shapes = list(prs.slides[0].shapes)
    assert len(shapes) == 1

    polygon = shapes[0]
    assert polygon.shape_type in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM)
    assert polygon.fill.fore_color.rgb == RGBColor(0xFF, 0xAA, 0x00)
    assert polygon.line.color.rgb == RGBColor(0x33, 0x33, 0x33)
