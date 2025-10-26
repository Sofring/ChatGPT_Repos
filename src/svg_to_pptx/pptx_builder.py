"""Generate PowerPoint presentations from the parsed SVG representation."""

from __future__ import annotations

from pathlib import Path
from typing import Sequence, Tuple

from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor

from .models import Circle, Document, Ellipse, Line, Path, Polyline, Rect, ShapeType, Text, TextStyle

PX_TO_EMU = 9525  # 1 px (assuming 96 DPI) -> EMU


def px_to_emu(value: float) -> Emu:
    return Emu(int(value * PX_TO_EMU))


def _set_fill(shape, style):
    fill_color = getattr(style, "fill_color", None)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor.from_string(fill_color.strip("#").upper())
    else:
        shape.fill.background()


def _set_line(shape, style):
    stroke_color = getattr(style, "stroke_color", None)
    if stroke_color:
        line = shape.line
        line.color.rgb = RGBColor.from_string(stroke_color.strip("#").upper())
        width = getattr(style, "stroke_width", 1.0)
        line.width = px_to_emu(width)
    else:
        shape.line.fill.background()


def _apply_style(shape, style):
    _set_fill(shape, style)
    _set_line(shape, style)


def _add_rect(slide, rect: Rect):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rect.rx or rect.ry else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type,
        px_to_emu(rect.x),
        px_to_emu(rect.y),
        px_to_emu(rect.width),
        px_to_emu(rect.height),
    )
    _apply_style(shape, rect.style)


def _add_circle(slide, circle: Circle):
    diameter = circle.r * 2
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        px_to_emu(circle.cx - circle.r),
        px_to_emu(circle.cy - circle.r),
        px_to_emu(diameter),
        px_to_emu(diameter),
    )
    _apply_style(shape, circle.style)


def _add_ellipse(slide, ellipse: Ellipse):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        px_to_emu(ellipse.cx - ellipse.rx),
        px_to_emu(ellipse.cy - ellipse.ry),
        px_to_emu(ellipse.rx * 2),
        px_to_emu(ellipse.ry * 2),
    )
    _apply_style(shape, ellipse.style)


def _add_line(slide, line: Line):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.LINE,
        px_to_emu(min(line.x1, line.x2)),
        px_to_emu(min(line.y1, line.y2)),
        px_to_emu(abs(line.x2 - line.x1)),
        px_to_emu(abs(line.y2 - line.y1)),
    )
    _apply_style(shape, line.style)


def _build_freeform(slide, points: Sequence[Tuple[float, float]], closed: bool):
    if not points or len(points) < 2:
        return None
    start = points[0]
    builder = slide.shapes.build_freeform(MSO_AUTO_SHAPE_TYPE.FREEFORM, px_to_emu(start[0]), px_to_emu(start[1]))
    for point in points[1:]:
        builder.add_line_segment(px_to_emu(point[0]), px_to_emu(point[1]))
    if closed:
        builder.add_line_segment(px_to_emu(start[0]), px_to_emu(start[1]))
    return builder.convert_to_shape()


def _add_polyline(slide, poly: Polyline):
    shape = _build_freeform(slide, poly.points, poly.closed)
    if shape:
        _apply_style(shape, poly.style)


def _add_path(slide, path: Path):
    for segment in path.segments:
        shape = _build_freeform(slide, segment.points, segment.closed)
        if shape:
            _apply_style(shape, path.style)


def _add_text(slide, text: Text):
    width = px_to_emu(200)
    height = px_to_emu(80)
    textbox = slide.shapes.add_textbox(px_to_emu(text.x), px_to_emu(text.y), width, height)
    frame = textbox.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text.text
    if text.style.fill_color:
        run.font.color.rgb = RGBColor.from_string(text.style.fill_color.strip("#").upper())
    if isinstance(text.style, TextStyle) and text.style.font_size:
        run.font.size = Pt(text.style.font_size * 0.75)
    if isinstance(text.style, TextStyle) and text.style.font_family:
        run.font.name = text.style.font_family


def build_presentation(document: Document, output_path: Path | str) -> Path:
    prs = Presentation()
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    slide = prs.slides.add_slide(blank_layout)

    if document.width:
        prs.slide_width = px_to_emu(document.width)
    if document.height:
        prs.slide_height = px_to_emu(document.height)

    for shape in document.shapes:
        if shape.shape_type == ShapeType.RECT:
            _add_rect(slide, shape)
        elif shape.shape_type == ShapeType.CIRCLE:
            _add_circle(slide, shape)
        elif shape.shape_type == ShapeType.ELLIPSE:
            _add_ellipse(slide, shape)
        elif shape.shape_type == ShapeType.LINE:
            _add_line(slide, shape)
        elif shape.shape_type in {ShapeType.POLYLINE, ShapeType.POLYGON}:
            _add_polyline(slide, shape)
        elif shape.shape_type == ShapeType.PATH:
            _add_path(slide, shape)
        elif shape.shape_type == ShapeType.TEXT:
            _add_text(slide, shape)

    output_path = Path(output_path)
    prs.save(output_path)
    return output_path

